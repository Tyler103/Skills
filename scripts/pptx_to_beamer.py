#!/usr/bin/env python3
"""
pptx_to_beamer.py  —  Convert a PowerPoint (.pptx) to a LaTeX Beamer document.

How it works (high level):
  1. Parse the .pptx file using python-pptx.
  2. For each slide, identify its title, bullet points, tables, images, and
     free-floating text boxes.
  3. Map each element to the closest LaTeX Beamer equivalent.
  4. Write a complete, compilable .tex file with a preamble, title frame, and
     one \\begin{frame}...\\end{frame} block per slide.

Usage:
    python pptx_to_beamer.py <input.pptx> [output.tex] [--theme THEME]

Options:
    --theme THEME   Override the inferred Beamer theme (e.g. Metropolis, Madrid).
                    Use "auto" (default) to let the script infer a theme from the
                    presentation's color palette.

Dependencies:
    pip install python-pptx pillow
"""

import sys
import os
import re
import argparse
from pathlib import Path

# Try to import python-pptx. If it isn't installed, print a helpful message and exit.
try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    from pptx.dml.color import RGBColor
    import pptx.shapes.picture
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx pillow")
    sys.exit(1)


# ──────────────────────────────────────────────────────────────────────────────
# LaTeX character escaping
#
# LaTeX treats these characters as special syntax: & % $ # _ { } ~ ^ \  < >
# If any of them appear in plain slide text they must be escaped, otherwise
# pdflatex will either throw an error or silently misrender the output.
#
# The table is ordered with backslash first because every other replacement
# introduces a backslash — processing backslash last would double-escape them.
# ──────────────────────────────────────────────────────────────────────────────

_ESCAPE_TABLE = [
    ("\\", r"\textbackslash{}"),  # must be first — other replacements add backslashes
    ("&",  r"\&"),                # column separator in tables
    ("%",  r"\%"),                # comment character
    ("$",  r"\$"),                # math mode delimiter
    ("#",  r"\#"),                # macro parameter character
    ("_",  r"\_"),                # subscript in math mode
    ("{",  r"\{"),                # grouping delimiter
    ("}",  r"\}"),                # grouping delimiter
    ("~",  r"\textasciitilde{}"), # non-breaking space
    ("^",  r"\textasciicircum{}"),# superscript in math mode
    ("<",  r"\textless{}"),       # less-than (not valid outside math without escaping)
    (">",  r"\textgreater{}"),    # greater-than
]

def escape_latex(text: str) -> str:
    """
    Replace every LaTeX-special character in `text` with its safe equivalent.

    Called on any raw string extracted from a pptx shape before it is written
    into the .tex output. Does NOT wrap the result in any environment — it only
    handles character-level escaping.
    """
    for char, replacement in _ESCAPE_TABLE:
        text = text.replace(char, replacement)
    return text


# ──────────────────────────────────────────────────────────────────────────────
# Rich-text run formatting
#
# Inside a PowerPoint text frame, text is stored in "runs" — contiguous spans
# that share the same font properties (bold, italic, underline, color, size).
# Each paragraph is made up of one or more runs.
#
# We iterate over runs rather than reading paragraph.text directly so we can
# honour per-run formatting and wrap each run in the appropriate LaTeX command.
# ──────────────────────────────────────────────────────────────────────────────

def format_run(run) -> str:
    """
    Convert a single text run to LaTeX, applying bold/italic/underline wrapping.

    Precedence rules:
      - Bold + italic together → \\textbf{\\textit{...}}
      - Bold alone             → \\textbf{...}
      - Italic alone           → \\textit{...}
      - Underline              → \\underline{...} applied on top of any above

    Whitespace-only runs are returned as-is (no wrapping needed).
    """
    text = escape_latex(run.text)

    # Skip runs that are purely whitespace — wrapping them would produce
    # \\textbf{ } which is harmless but noisy.
    if not text.strip():
        return text

    bold      = run.font.bold
    italic    = run.font.italic
    underline = run.font.underline

    # Apply bold/italic first, then underline on the outside.
    if bold and italic:
        text = rf"\textbf{{\textit{{{text}}}}}"
    elif bold:
        text = rf"\textbf{{{text}}}"
    elif italic:
        text = rf"\textit{{{text}}}"

    # \\underline wraps whatever bold/italic combination was already applied.
    if underline:
        text = rf"\underline{{{text}}}"

    return text


def paragraph_to_latex(para) -> str:
    """
    Concatenate all runs in a paragraph into a single LaTeX string.

    Each run is passed through format_run() so rich formatting is preserved.
    The paragraph-level text (para.text) is NOT used directly because it strips
    run boundaries and loses all font information.
    """
    return "".join(format_run(r) for r in para.runs)


# ──────────────────────────────────────────────────────────────────────────────
# Shape-type helpers
#
# PowerPoint slides contain "shapes". Each shape has a type (picture, table,
# text box, etc.) and may also have a "placeholder format" that describes its
# role on the slide (title, body/bullets, subtitle, date, footer, ...).
#
# We care about three categories:
#   1. Title placeholders  → become the frame title: \begin{frame}{Title}
#   2. Body placeholders   → contain the bullet-point content
#   3. Everything else     → free-floating text boxes, wrapped in \begin{block}
# ──────────────────────────────────────────────────────────────────────────────

# The set of placeholder types that represent a slide title.
# CENTER_TITLE is used on title slides; SUBTITLE is treated as a title too
# because it usually contains the main heading on the first slide.
_TITLE_PH_TYPES = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.SUBTITLE,
}

def _ph_type(shape):
    """
    Safely return the PP_PLACEHOLDER enum value for a shape, or None.

    Some shapes have corrupt or missing placeholder data — the try/except
    prevents those from crashing the whole conversion.
    """
    try:
        ph = shape.placeholder_format
        return ph.type if ph else None
    except (ValueError, AttributeError):
        return None


def is_title(shape) -> bool:
    """True if the shape is a title or center-title placeholder."""
    return shape.has_text_frame and _ph_type(shape) in _TITLE_PH_TYPES


def is_body(shape) -> bool:
    """
    True for body/content placeholders that hold bullet-point text.

    Any placeholder that is NOT a title type is considered body content.
    This includes BODY, OBJECT, and other content placeholder variants.
    """
    t = _ph_type(shape)
    # t must be a real placeholder type (not None) and not one of the title types.
    return t is not None and t not in _TITLE_PH_TYPES


# ──────────────────────────────────────────────────────────────────────────────
# Table conversion
#
# pptx tables are accessed via shape.table and are structured as rows of cells.
# We convert them to a LaTeX tabular environment with vertical and horizontal
# rules. The column spec uses left-aligned (l) columns — one per table column.
# ──────────────────────────────────────────────────────────────────────────────

def table_to_latex(table) -> str:
    """
    Convert a python-pptx Table object into a LaTeX tabular environment.

    Layout decisions:
      - All columns are left-aligned (l). To right-align numbers, edit the
        col_spec string after conversion (e.g. replace 'l' with 'r' as needed).
      - Every row gets a \\hline above and below it for a simple grid style.
      - Cell text is escaped for LaTeX special characters.

    Returns an empty string if the table has no rows.
    """
    rows = table.rows
    if not rows:
        return ""  # nothing to render

    # Build the column specification string, e.g. "|l|l|l|" for 3 columns.
    num_cols = len(rows[0].cells)
    col_spec = "|" + "|".join(["l"] * num_cols) + "|"

    lines = [
        r"\begin{center}",          # centre the table on the slide
        rf"\begin{{tabular}}{{{col_spec}}}",
        r"\hline",                  # top border
    ]

    for row in enumerate(rows):
        cells = []
        for cell in row[1].cells:
            # cell.text_frame holds the text; fall back to empty string if missing
            raw = cell.text_frame.text.strip() if cell.text_frame else ""
            cells.append(escape_latex(raw))

        # LaTeX table rows: cells joined with & and terminated with \\
        row_str = " & ".join(cells) + r" \\"
        lines.append(row_str)
        lines.append(r"\hline")    # horizontal rule after every row

    lines += [r"\end{tabular}", r"\end{center}"]
    return "\n".join(lines)


# ──────────────────────────────────────────────────────────────────────────────
# Bullet-list builder
#
# PowerPoint bullet paragraphs have a `level` attribute (0 = top level,
# 1 = first indent, 2 = second indent, ...).  We map these to nested
# \begin{itemize} environments.
#
# The algorithm uses a stack (open_envs) to track which environments are
# currently open, and opens/closes them as the indentation level changes.
# ──────────────────────────────────────────────────────────────────────────────

def bullets_to_latex(items: list[tuple[int, str]]) -> list[str]:
    """
    Convert a list of (indent_level, latex_text) tuples into nested itemize blocks.

    Example input:
        [(0, "First bullet"), (1, "Sub-bullet"), (0, "Second bullet")]

    Example output:
        \\begin{itemize}
          \\item First bullet
          \\begin{itemize}
            \\item Sub-bullet
          \\end{itemize}
          \\item Second bullet
        \\end{itemize}

    Args:
        items: List of (level, text) pairs. Level 0 = outermost bullet.

    Returns:
        List of LaTeX lines (not yet joined with newlines).
    """
    lines = []
    depth = -1          # current nesting depth (-1 means no env is open yet)
    open_envs = []      # stack of open itemize depths, used to close them in order

    for level, text in items:
        # Close environments that are deeper than the current bullet's level.
        while depth > level:
            lines.append("  " * depth + r"\end{itemize}")
            open_envs.pop()
            depth -= 1

        # Open new environments until we reach the current bullet's depth.
        while depth < level:
            depth += 1
            lines.append("  " * depth + r"\begin{itemize}")
            open_envs.append(depth)

        # Emit the bullet item, indented to match the current depth.
        lines.append("  " * (depth + 1) + rf"\item {text}")

    # After all bullets are processed, close any environments still on the stack.
    while open_envs:
        d = open_envs.pop()
        lines.append("  " * d + r"\end{itemize}")

    return lines


# ──────────────────────────────────────────────────────────────────────────────
# Single-slide conversion
#
# Each slide becomes one \begin{frame}...\end{frame} block. We make two passes
# over the slide's shapes:
#   Pass 1 (implicit — for loop): classify and collect each shape's content.
#   Pass 2: assemble the collected content into ordered LaTeX lines.
#
# Output order within a frame: bullets → tables → free text boxes → images.
# ──────────────────────────────────────────────────────────────────────────────

def extract_slide(slide, images_dir: str, slide_num: int) -> str:
    """
    Convert a single pptx Slide object to a LaTeX Beamer frame string.

    Args:
        slide:      A python-pptx Slide object.
        images_dir: Absolute path to the directory where extracted images are saved.
        slide_num:  1-based slide number, used to name extracted image files.

    Returns:
        A multi-line string containing the complete LaTeX frame, from
        \\begin{frame}{...} to \\end{frame}.
    """
    title_text = ""                         # the slide's title string
    bullet_items: list[tuple[int, str]] = []# (indent_level, latex_text) pairs
    table_blocks: list[str] = []            # rendered tabular environments
    freetext_blocks: list[str] = []         # text from non-placeholder text boxes
    image_names: list[str] = []             # filenames of extracted images

    # ── Iterate over every shape on the slide ────────────────────────────────
    for shape in slide.shapes:

        # ── Images ──────────────────────────────────────────────────────────
        # pptx.shapes.picture.Picture covers raster images (PNG, JPEG, etc.)
        # embedded directly in the slide. We save the raw binary blob to disk
        # and record the filename for later use in \includegraphics.
        if isinstance(shape, pptx.shapes.picture.Picture):
            try:
                # Name images as slide<N>_img<M>.png for easy identification.
                img_name = f"slide{slide_num}_img{len(image_names)+1}.png"
                img_path = os.path.join(images_dir, img_name)
                with open(img_path, "wb") as fh:
                    fh.write(shape.image.blob)  # raw image bytes from the pptx
                image_names.append(img_name)
            except Exception as e:
                # Non-fatal: warn and skip rather than aborting the whole conversion.
                print(f"  [warn] Could not extract image on slide {slide_num}: {e}")
            continue  # move to the next shape — pictures have no text frame

        # ── Tables ──────────────────────────────────────────────────────────
        # MSO_SHAPE_TYPE.TABLE identifies table shapes. The table object holds
        # a grid of rows and cells that we convert to a LaTeX tabular.
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            try:
                table_blocks.append(table_to_latex(shape.table))
            except Exception as e:
                print(f"  [warn] Could not convert table on slide {slide_num}: {e}")
            continue

        # ── All other shapes must have a text frame to be useful ─────────────
        if not shape.has_text_frame:
            # Connectors, decorative shapes, SmartArt, charts, etc. are silently
            # skipped — there is no reliable way to represent them in LaTeX.
            continue

        # ── Title placeholder ────────────────────────────────────────────────
        # The title becomes the argument to \begin{frame}{...}.
        # We read the whole text frame text (not per-run) because we only need
        # the plain string here — formatting in titles is unusual.
        if is_title(shape):
            title_text = escape_latex(shape.text_frame.text.strip())
            continue

        # ── Body placeholder (bullet points) ────────────────────────────────
        # Body placeholders hold the main bullet-point content of a slide.
        # Each paragraph maps to one bullet item; para.level gives its depth.
        if is_body(shape):
            for para in shape.text_frame.paragraphs:
                # paragraph_to_latex preserves bold/italic/underline per run.
                text = paragraph_to_latex(para).strip()
                if text:  # skip empty paragraphs (blank lines in the pptx)
                    bullet_items.append((para.level, text))

        else:
            # ── Free-floating text box ───────────────────────────────────────
            # Non-placeholder text boxes (e.g. manually inserted text boxes,
            # speaker-note callouts) are collected as plain blocks and wrapped
            # in \begin{block}{} so they appear as a shaded Beamer block.
            lines = []
            for para in shape.text_frame.paragraphs:
                text = paragraph_to_latex(para).strip()
                if text:
                    lines.append(text)
            if lines:
                freetext_blocks.append("\n".join(lines))

    # ── Assemble the frame LaTeX ─────────────────────────────────────────────
    # Build the list of output lines, then join them at the end.
    frame_lines: list[str] = []

    # Frame opening: use the slide title as the frame title, or empty braces
    # if no title was found (Beamer still renders the frame, just without a bar).
    title_arg = f"{{{title_text}}}" if title_text else "{}"
    frame_lines.append(rf"\begin{{frame}}{title_arg}")

    # 1. Bullet points (nested itemize environments)
    if bullet_items:
        frame_lines.extend(bullets_to_latex(bullet_items))

    # 2. Tables (each rendered as a centred tabular)
    for tbl in table_blocks:
        frame_lines.append("")   # blank line for readability
        frame_lines.append(tbl)

    # 3. Free-floating text boxes (wrapped in Beamer's coloured block environment)
    for block in freetext_blocks:
        frame_lines.append("")
        frame_lines.append(r"\begin{block}{}")
        frame_lines.append(block)
        frame_lines.append(r"\end{block}")

    # 4. Images (placed below all text content, scaled to 80% of the text width)
    for img in image_names:
        # Path is relative to the .tex file location so pdflatex can find it
        # when compiled from the same directory.
        frame_lines.append(rf"  \includegraphics[width=0.8\textwidth]{{beamer_images/{img}}}")

    frame_lines.append(r"\end{frame}")
    return "\n".join(frame_lines)


# ──────────────────────────────────────────────────────────────────────────────
# Theme inference
#
# A good Beamer theme makes the presentation look intentional. Rather than
# always defaulting to Madrid, we try two heuristics:
#
#   1. Match the pptx theme name (embedded in the XML) to a known Beamer theme.
#      Many pptx files use Microsoft-named themes (Ion, Facet, Circuit, etc.)
#      that have rough Beamer equivalents in style or colour.
#
#   2. Analyse the dominant background colour:
#      - Very dark → Metropolis (modern dark theme)
#      - Blue-heavy → Warsaw
#      - Red/maroon → AnnArbor
#      - Green-heavy → Hannover
#
#   3. Default to Madrid — a widely recognised professional theme.
# ──────────────────────────────────────────────────────────────────────────────

# Mapping from substrings of pptx theme names (lowercase) to Beamer theme names.
# Extend this dict to support more pptx-to-Beamer mappings.
_THEME_NAME_MAP = {
    "metropolitan": "Metropolis",
    "metro":        "Metropolis",
    "office":       "Madrid",
    "circuit":      "Berlin",
    "facet":        "AnnArbor",
    "ion":          "Warsaw",
    "wood":         "Hannover",
    "slate":        "CambridgeUS",
    "retrospect":   "Boadilla",
    "wisp":         "Luebeck",
    "organic":      "Bergen",
}

def _get_theme_name_from_pptx(prs: Presentation) -> str | None:
    """
    Read the theme name from the pptx XML and return it in lowercase.

    The theme element lives in the slide master's XML. Its `name` attribute
    usually contains a descriptive string like "Office Theme" or "Ion Boardroom".
    Returns None if the name cannot be read (e.g. corrupt file, no theme).
    """
    try:
        theme_elem = prs.slide_masters[0].theme.element
        name = theme_elem.get("name", "")
        return name.lower() if name else None
    except Exception:
        return None


def _dominant_color(prs: Presentation) -> tuple[int, int, int] | None:
    """
    Return the (R, G, B) tuple of the first slide's solid background colour.

    Tries the slide background first, then falls back to the slide master
    background (which most slides inherit if they don't override it).
    Returns None if no solid background colour can be determined.
    """
    # Try the first slide's background fill
    try:
        bg_fill = prs.slides[0].background.fill
        if bg_fill.type is not None:
            clr = bg_fill.fore_color.rgb   # hex string like "1F3864"
            return int(clr[0:2], 16), int(clr[2:4], 16), int(clr[4:6], 16)
    except Exception:
        pass

    # Fall back to the slide master's background
    try:
        bg_fill = prs.slide_masters[0].background.fill
        if bg_fill.type is not None:
            clr = bg_fill.fore_color.rgb
            return int(clr[0:2], 16), int(clr[2:4], 16), int(clr[4:6], 16)
    except Exception:
        pass

    return None


def infer_beamer_theme(prs: Presentation) -> str:
    """
    Heuristically choose a Beamer theme that suits the presentation's style.

    Strategy:
      1. Try to match the embedded pptx theme name against _THEME_NAME_MAP.
      2. Analyse the dominant background colour using the ITU-R BT.601
         perceptual brightness formula: Y = 0.299R + 0.587G + 0.114B.
      3. Default to Madrid if nothing better can be determined.
    """
    # Step 1: name-based matching
    pptx_name = _get_theme_name_from_pptx(prs)
    if pptx_name:
        for keyword, beamer in _THEME_NAME_MAP.items():
            if keyword in pptx_name:
                return beamer

    # Step 2: colour-based matching
    color = _dominant_color(prs)
    if color:
        r, g, b = color
        # Perceptual brightness (0 = black, 255 = white)
        brightness = (r * 299 + g * 587 + b * 114) // 1000

        if brightness < 60:
            return "Metropolis"   # very dark background → modern dark theme
        if r < 80 and b > 120:
            return "Warsaw"       # blue-dominant background
        if r > 140 and g < 80:
            return "AnnArbor"     # red or maroon dominant
        if g > 120 and r < 100:
            return "Hannover"     # green dominant

    # Step 3: fall back to a universally recognised professional default
    return "Madrid"


# ──────────────────────────────────────────────────────────────────────────────
# Presentation title extraction
# ──────────────────────────────────────────────────────────────────────────────

def get_presentation_title(prs: Presentation) -> str:
    """
    Extract a title string from the first slide to use in the LaTeX \\title{}.

    Reads the first title placeholder found on slide 1. If no title is found
    (e.g. the first slide is a blank or image-only slide), returns "Presentation"
    as a safe fallback.
    """
    if not prs.slides:
        return "Presentation"
    for shape in prs.slides[0].shapes:
        if is_title(shape) and shape.text_frame.text.strip():
            return escape_latex(shape.text_frame.text.strip())
    return "Presentation"


# ──────────────────────────────────────────────────────────────────────────────
# Document assembly
# ──────────────────────────────────────────────────────────────────────────────

def build_preamble(title: str, theme: str) -> str:
    """
    Return the LaTeX document preamble as a string.

    Includes:
      - \\documentclass{beamer}
      - The chosen Beamer theme
      - Standard packages: graphicx (images), hyperref (links), amsmath (math),
        booktabs (publication-quality tables), inputenc (UTF-8 support)
      - \\title and \\date metadata
      - \\begin{document} and an auto-generated title frame

    Args:
        title: The presentation title (already LaTeX-escaped).
        theme: The Beamer theme name (e.g. "Madrid", "Metropolis").
    """
    return rf"""% Generated by pptx_to_beamer.py
\documentclass{{beamer}}
\usetheme{{{theme}}}
\usepackage{{graphicx}}      % required for \includegraphics
\usepackage{{hyperref}}      % clickable URLs and cross-references
\usepackage{{amsmath}}       % math environments and symbols
\usepackage{{booktabs}}      % \toprule / \midrule / \bottomrule for tables
\usepackage[utf8]{{inputenc}}% allow UTF-8 characters in source

\title{{{title}}}
\date{{\today}}

\begin{{document}}

% Auto-generated title slide
\begin{{frame}}
  \titlepage
\end{{frame}}

"""


def convert(input_path: str, output_path: str, theme_override: str = "auto") -> None:
    """
    Main conversion function: read a .pptx and write a .tex file.

    Steps:
      1. Open the pptx with python-pptx.
      2. Create the beamer_images/ output directory for extracted images.
      3. Determine the Beamer theme (inferred or overridden).
      4. Convert each slide to a LaTeX frame string.
      5. Write the preamble + all frames + \\end{document} to the output file.
      6. Clean up the images directory if no images were extracted.

    Args:
        input_path:     Path to the source .pptx file.
        output_path:    Path to write the output .tex file.
        theme_override: Beamer theme name, or "auto" to infer from the pptx.
    """
    prs = Presentation(input_path)

    # Create a beamer_images/ folder next to the output .tex file.
    # Images are stored here and referenced with relative paths so the .tex
    # file is portable — it compiles correctly from any machine as long as
    # beamer_images/ is present in the same directory.
    out_dir = os.path.dirname(os.path.abspath(output_path))
    images_dir = os.path.join(out_dir, "beamer_images")
    os.makedirs(images_dir, exist_ok=True)

    doc_title = get_presentation_title(prs)

    # Resolve the theme: either infer it from the pptx or use the user's choice.
    if theme_override.lower() == "auto":
        theme = infer_beamer_theme(prs)
        print(f"  Inferred Beamer theme: {theme}")
    else:
        theme = theme_override
        print(f"  Using Beamer theme: {theme}")

    # Convert each slide, printing progress to stdout.
    frames = []
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        print(f"  Converting slide {i}/{total} …", end="\r")
        frames.append(extract_slide(slide, images_dir, i))

    print()  # move to a new line after the \r progress indicator

    # Write the complete .tex file.
    with open(output_path, "w", encoding="utf-8") as fh:
        fh.write(build_preamble(doc_title, theme))
        fh.write("\n\n".join(frames))          # blank line between each frame
        fh.write("\n\n\\end{document}\n")

    print(f"Done: {output_path}")

    # Report on extracted images (or clean up the empty directory).
    extracted = os.listdir(images_dir)
    if extracted:
        print(f"Images ({len(extracted)}): {images_dir}/")
    else:
        os.rmdir(images_dir)   # no images were found — remove the empty folder

    print(f"\nTo compile:\n  pdflatex \"{os.path.basename(output_path)}\"")


# ──────────────────────────────────────────────────────────────────────────────
# CLI entry point
# ──────────────────────────────────────────────────────────────────────────────

def main():
    """
    Parse command-line arguments and kick off the conversion.

    Argument summary:
      input   (required) — path to the source .pptx file
      output  (optional) — path for the output .tex file;
                           defaults to <input_stem>_beamer.tex in the cwd
      --theme (optional) — Beamer theme name or "auto" (default)
    """
    parser = argparse.ArgumentParser(
        description="Convert a .pptx file to a LaTeX Beamer presentation."
    )
    parser.add_argument("input",  help="Input .pptx file")
    parser.add_argument("output", nargs="?", help="Output .tex file (optional)")
    parser.add_argument(
        "--theme",
        default="auto",
        help='Beamer theme name, or "auto" to infer from the presentation (default: auto)',
    )
    args = parser.parse_args()

    inp = args.input
    # Default output filename: strip the .pptx extension and add _beamer.tex
    out = args.output or (Path(inp).stem + "_beamer.tex")
    convert(inp, out, theme_override=args.theme)


# Run main() only when the script is executed directly, not when it is imported
# as a module (e.g. during testing).
if __name__ == "__main__":
    main()
